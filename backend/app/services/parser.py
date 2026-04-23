import pandas as pd
from docx import Document
from pptx import Presentation
import json
import os
import pytesseract
from pdf2image import convert_from_path
import pdfplumber
import platform

# main function: parse 3 types of file
# converts tabular data into json or list of dictionary
def extract_from_excel(file_path):
    try:
        # assume the first row is the header
        df = pd.read_excel(file_path)
        # drop the rows that are completely empty 
        df = df.dropna(how='all')
        # Convert to a list of dictionaries (each row beacomes one record)
        records = df.to_dict(orient="records")
        return json.dumps(records, ensure_ascii=False, indent=2)
    except Exception as e:
        return f"Excel parsing failed: {e}"

# extracts plain text paragraph by paragraph
def extract_from_word(file_path):
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text.strip())
        return "\n".join(full_text)
    except Exception as e:
        return f"Word parsing failed: {e}"
    
# iterates through every text box on every slide
def extract_from_ppt(file_path):
    try:
        prs = Presentation(file_path)
        text_runs = []

        for slide_num, slide in enumerate(prs.slides):
            slide_text = f"--- Slide {slide_num + 1} ---"
            text_runs.append(slide_text)
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
                    
                elif shape.has_table:
                    text_runs.append("  [tabular data detected]:")
                    for row in shape.table.rows:
                        row_data = [cell.text_frame.text.replace('\n', ' ').strip() for cell in row.cells]
                        text_runs.append("  | " + " | ".join(row_data) + " |")
            
        return "\n".join(text_runs)
        
    except Exception as e:
        return f"PPT parsing failed: {e}"


# extract pdf text page by page
def extract_text_pdf(file_path):
    try:
        full_text = []

        # open pdf file
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                full_text.append(f"\n --- Page {page_num + 1} ---")

                text = page.extract_text()

                if text:
                    full_text.append(text)
        return "\n".join(full_text)
    
    except Exception as e:
        return f"PDF parsing failed: {e}"

if platform.system() == "Windows"
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    poppler_bin_path = r'C:\Users\YIXIN\Downloads\Release-25.12.0-0\poppler-25.12.0\Library\bin'
else:
    poppler_bin_path = None

def extract_with_local_ocr(file_path):
    try:
        # convert pdf into picture
        images =convert_from_path(file_path, dpi=300, poppler_path=poppler_bin_path)

        full_text = []

        for i,image in enumerate(images):
            print (f"[OCR] detecting page {i + 1}")
            text = pytesseract.image_to_string(image, lang = 'eng')

            if text.strip():
                full_text.append(f"\n --- Page {i + 1}---")
                full_text.append(text)

        print("[OCR] parsing completed")
        return "\n".join(full_text)
            
    except Exception as e:
        return f"OCR PDF parsing failed: {e}"
    
# decides whether to use text extraction or OCR
def extract_from_pdf(file_path):
    try:
        with pdfplumber.open(file_path) as pdf:
            requires_ocr = False

            if len(pdf.pages) == 0:
                return "Error: Empty PDF."
            
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                images = page.images

                # text only
                if len (images) > 0:
                    requires_ocr = True
                    break
                
            if not requires_ocr:
                return extract_text_pdf(file_path)
            
            else:
                return extract_with_local_ocr(file_path)
        
    except Exception as e:
        return f"PDF parsing failed: {e}"