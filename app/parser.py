import pandas as pd
from docx import Document
from pptx import Presentation
import json
import pdfplumber

# main function: parse 3 types of file

# converts tabular data into json or list of dictionary
def extract_from_excel(file_path):
    try:
        # asume the first row is the header
        df = pd.read_excel(file_path)
        # drop the rows that are completely empty 
        df = df.dropna(how='all')
        # Convert to a ist of dictionaries (each row beacomes one record)
        records = df.to_dict(orient="records")
        return json.dumps(records, ensure_ascii=False, indent=2)
    except Exception as e:
        return f"Excel parsing failed: {e}"

# extracts plain text paragrapg by paragraph
def extract_from_word(file_path):
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text.strip())
        return "\n".join(full_text)
    except Exception as e:
        return f"Word parsing failed: {e}"
    
# iterates through every text box on every slide
def extract_from_ppt(file_path):
    try:
        prs = Presentation(file_path)
        text_runs = []

        for slide_num, slide in enumerate(prs.slides):
            slide_text = f"--- Slide {slide_num + 1} ---"
            text_runs.append(slide_text)
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
                    
                elif shape.has_table:
                    text_runs.append("  [tabular data detected]:")
                    for row in shape.table.rows:
                        row_data = [cell.text_frame.text.replace('\n', ' ').strip() for cell in row.cells]
                        text_runs.append("  | " + " | ".join(row_data) + " |")
            
        return "\n".join(text_runs)
        
    except Exception as e:
        return f"PPT parsing failed: {e}"

'''
# extract pdf text page by page
def extract_from_pdf(file_path):
    try:
        full_text = []

        # open pdf file
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                full_text.append(f"\n --- Page {page_num + 1} ---")

                text = page.extract_text()

                if text:
                    full_text.append(text)

        return "\n".join(full_text)
    
    except Exception as e:
        return f"PDF parsing failed: {e}"
'''

# Testing
if __name__ == "__main__":
    print("--- Testing Excel Extraction ---")
    excel_data = extract_from_excel(r"C:\Users\YIXIN\Downloads\Lampiran A2 Senarai pelajar FAIX.xlsx")
    print(excel_data)

    print("\n--- Testing Word Extraction ---")
    word_text = extract_from_word(r"C:\Users\YIXIN\Downloads\EPI-script.docx")
    print(word_text)

    print("\n--- Testing PPT Extraction ---")
    ppt_text = extract_from_ppt(r"C:\Users\YIXIN\Downloads\Chapter 6_Estimation_v4_sakinah.pptx")
    print(ppt_text)

    #print("\n--- Testing PDF Extraction ---")
    #pdf_text = extract_from_pdf(r"C:\Users\YIXIN\Downloads\4. KBS Life Cycle SEM120252026.pdf")